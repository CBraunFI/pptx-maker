import io, base64
import os
import unicodedata
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from typing import List, Union

def hex_to_rgb(hexstr: str):
    hexstr = hexstr.lstrip("#")
    return RGBColor(int(hexstr[0:2],16), int(hexstr[2:4],16), int(hexstr[4:6],16))

def sanitize_text(text: str) -> str:
    if text is None:
        return ""
    # Replace common problematic Unicode characters
    replacements = {
        '\u2013': '-',  # En dash
        '\u2014': '-',  # Em dash
        '\u2018': "'",  # Left single quotation mark
        '\u2019': "'",  # Right single quotation mark
        '\u201C': '"',  # Left double quotation mark
        '\u201D': '"',  # Right double quotation mark
        '\u2026': '...',  # Horizontal ellipsis
        '\u00A0': ' ',  # Non-breaking space
    }
    for u, a in replacements.items():
        text = text.replace(u, a)
    text = unicodedata.normalize('NFKD', text)
    # keep latin-1 safe (avoid odd glyphs)
    return text.encode('ascii', 'ignore').decode('ascii')

def add_logos(slide, synk_logo_path, client_logo_path, slide_width, slide_height):
    # SYNK logo - bottom right corner
    if synk_logo_path and os.path.exists(synk_logo_path):
        logo_height = Inches(0.4)
        logo_width = Inches(1.2)
        left = slide_width - logo_width - Inches(0.3)
        top = slide_height - logo_height - Inches(0.2)
        try:
            slide.shapes.add_picture(synk_logo_path, left, top, height=logo_height)
        except Exception:
            pass
    # Client logo - bottom left corner
    if client_logo_path and os.path.exists(client_logo_path):
        logo_height = Inches(0.4)
        left = Inches(0.3)
        top = slide_height - logo_height - Inches(0.2)
        try:
            slide.shapes.add_picture(client_logo_path, left, top, height=logo_height)
        except Exception:
            pass

def add_title_slide(prs, meta, slide):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # background (primary)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(meta["style"]["colors"]["primary"])
    # title
    tb = s.shapes.add_textbox(Inches(1), Inches(1.7), prs.slide_width - Inches(2), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = sanitize_text(slide.get("title",""))
    p.font.name = "Arial"; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)
    tb.text_frame.word_wrap = True
    # subtitle
    sub = s.shapes.add_textbox(Inches(1), Inches(2.7), prs.slide_width - Inches(2), Inches(0.8))
    sp = sub.text_frame.paragraphs[0]
    sp.text = sanitize_text(slide.get("subtitle") or (meta.get("deckSubtitle") or ""))
    sp.font.name = "Arial"; sp.font.size = Pt(20); sp.font.color.rgb = RGBColor(255,255,255)
    sub.text_frame.word_wrap = True
    return s

def _normalize_content_from_slide(slide: dict) -> List[str]:
    """
    Builds a flat list of strings to render as paragraphs from various schema variants.
    Priority:
      - 'content' if present (string or list)
      - else 'text' (+ optional 'bullets' or 'items')
      - else 'items' / 'bullets'
      - for team: 'members'/'trainers' dicts → lines
      - for contact: 'contact' dict → lines
    """
    out: List[str] = []

    # Explicit content wins
    if "content" in slide and slide["content"] not in (None, "", []):
        c = slide["content"]
        if isinstance(c, str):
            out.append(c)
        elif isinstance(c, list):
            out.extend([str(x) for x in c if x is not None])
        return [sanitize_text(x) for x in out]

    # Contact block
    if "contact" in slide and isinstance(slide["contact"], dict):
        c = slide["contact"]
        line1 = " - ".join([x for x in [c.get("name"), c.get("role")] if x])
        line2 = c.get("email")
        line3 = c.get("phone")
        for ln in [line1, line2, line3]:
            if ln:
                out.append(ln)

    # Team members / trainers
    members = []
    if isinstance(slide.get("members"), list):
        members = slide["members"]
    elif isinstance(slide.get("trainers"), list):
        members = slide["trainers"]
    if members:
        for m in members:
            if isinstance(m, dict):
                name = m.get("name","")
                role = m.get("role")
                focus = m.get("focus")
                parts = [name]
                if role:  parts.append(role)
                if focus: parts.append(f"({focus})")
                out.append(" – ".join([p for p in parts if p]))
            else:
                out.append(str(m))

    # Base text + bullets/items
    base_text = slide.get("text")
    if isinstance(base_text, str) and base_text.strip():
        out.insert(0, base_text)  # lead paragraph first

    for key in ("bullets", "items"):
        lst = slide.get(key)
        if isinstance(lst, list) and lst:
            for x in lst:
                if isinstance(x, str):
                    out.append(x)
                elif isinstance(x, dict):
                    # e.g., investment items with label/value/note
                    label = x.get("label")
                    value = x.get("value")
                    note  = x.get("note")
                    if label or value or note:
                        triple = " – ".join([t for t in [label, value, note] if t])
                        out.append(triple)

    out = [sanitize_text(x) for x in out if x is not None]
    return out

def add_text_slide(prs, meta, slide, header="", synk_logo_path=None, client_logo_path=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_logos(s, synk_logo_path, client_logo_path, prs.slide_width, prs.slide_height)

    # header
    hdr = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(0.6))
    hp = hdr.text_frame.paragraphs[0]
    title_text = slide.get("title","")
    if header:
        title_text = f"{header} - {title_text}" if title_text else header
    hp.text = sanitize_text(title_text)
    hp.font.name = "Arial"; hp.font.size = Pt(28); hp.font.bold = True
    hp.font.color.rgb = hex_to_rgb(meta["style"]["colors"]["text"])
    hdr.text_frame.word_wrap = True

    # body
    left, top, width, height = Inches(0.5), Inches(1.3), Inches(9.0), Inches(3.8)
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True

    content_list = _normalize_content_from_slide(slide)
    if not content_list:
        tf.paragraphs[0].text = ""
        return s

    for i, item in enumerate(content_list):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Arial"; p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(meta["style"]["colors"]["text"])
    return s

def add_table_slide(prs, meta, slide, headers: List[str], rows: List[List[str]], synk_logo_path=None, client_logo_path=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_logos(s, synk_logo_path, client_logo_path, prs.slide_width, prs.slide_height)

    # header
    hdr = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(0.6))
    hp = hdr.text_frame.paragraphs[0]
    hp.text = sanitize_text(slide.get("title",""))
    hp.font.name = "Arial"; hp.font.size = Pt(28); hp.font.bold = True
    hp.font.color.rgb = hex_to_rgb(meta["style"]["colors"]["text"])

    # table
    table = s.shapes.add_table(len(rows)+1, len(headers), Inches(0.5), Inches(1.3), Inches(9.0), Inches(3.8)).table
    accent = hex_to_rgb(meta["style"]["colors"]["accent1"])
    for j,h in enumerate(headers):
        cell = table.cell(0,j)
        cell.text = sanitize_text(h)
        cell.fill.solid(); cell.fill.fore_color.rgb = accent
        cell.text_frame.paragraphs[0].font.name = "Arial"
        cell.text_frame.paragraphs[0].font.bold = True
    for i,r in enumerate(rows, start=1):
        for j,val in enumerate(r):
            cell = table.cell(i,j)
            cell.text = sanitize_text(val if val is not None else "")
            cell.text_frame.paragraphs[0].font.name = "Arial"
    return s

def build_pptx(deck: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    meta = deck["meta"]

    # logos
    synk_logo_path = meta.get("style", {}).get("logo")
    if synk_logo_path and not os.path.exists(synk_logo_path):
        for ext in ['.png', '.PNG', '.jpg', '.JPG', '.jpeg', '.JPEG']:
            base_name = os.path.splitext(synk_logo_path)[0]
            test_path = base_name + ext
            if os.path.exists(test_path):
                synk_logo_path = test_path
                break

    client_logo_path = meta.get("style", {}).get("clientLogo")
    if client_logo_path and not os.path.exists(client_logo_path):
        for ext in ['.png', '.PNG', '.jpg', '.JPG', '.jpeg', '.JPEG']:
            base_name = os.path.splitext(client_logo_path)[0]
            test_path = base_name + ext
            if os.path.exists(test_path):
                client_logo_path = test_path
                break

    # Slides
    for sl in deck["slides"]:
        t = sl.get("type","")
        if t == "title":
            add_title_slide(prs, meta, sl)

        elif t == "agenda":
            # Map 'items' → content
            sl2 = dict(sl)
            if "content" not in sl2:
                sl2["content"] = sl.get("items") or sl.get("bullets") or []
            add_text_slide(prs, meta, sl2, synk_logo_path=synk_logo_path, client_logo_path=client_logo_path)

        elif t in ["context","need","understanding","vision","approach","principles",
                   "architecture","transfer","digital","coaching","target_group","impact",
                   "about_synk","references","expertise","partners","next_steps","contact"]:
            # Normalize text/bullets/items/contact/members → content
            add_text_slide(prs, meta, sl, synk_logo_path=synk_logo_path, client_logo_path=client_logo_path)

        elif t == "modules_overview":
            headers = ["Modul","Dauer","Fokus"]
            rows = []
            for m in (sl.get("modules") or []):
                rows.append([m.get("title",""), m.get("duration",""), m.get("focus","")])
            add_table_slide(prs, meta, sl, headers, rows or [["—","—","—"]], synk_logo_path=synk_logo_path, client_logo_path=client_logo_path)

        elif t == "module_detail":
            add_text_slide(prs, meta, sl, header="Modul", synk_logo_path=synk_logo_path, client_logo_path=client_logo_path)

        elif t == "team":
            # Accept 'members' or 'trainers'
            sl2 = dict(sl)
            lines = _normalize_content_from_slide(sl)  # will pick up members/trainers
            if lines:
                sl2["content"] = lines
            add_text_slide(prs, meta, sl2, header="Team", synk_logo_path=synk_logo_path, client_logo_path=client_logo_path)

        elif t == "investment":
            # Prefer structured items [{label,value,note}], fallback to 'content' parsing
            items = sl.get("items")
            if isinstance(items, list) and items and isinstance(items[0], dict):
                headers = ["Position","Hinweis","Preis"]
                rows = []
                for it in items:
                    rows.append([
                        it.get("label",""),
                        it.get("note",""),
                        it.get("value","")
                    ])
                add_table_slide(prs, meta, sl, headers, rows, synk_logo_path=synk_logo_path, client_logo_path=client_logo_path)
            else:
                # legacy 'content' lines: "Label – Preis"
                headers = ["Position","Hinweis","Preis"]
                rows = []
                content = sl.get("content") or []
                for c in content:
                    if isinstance(c, str) and "–" in c:
                        left,right = c.split("–",1)
                        rows.append([left.strip(), "", right.strip()])
                    else:
                        rows.append([str(c), "", ""])
                add_table_slide(prs, meta, sl, headers, rows or [["—","","—"]], synk_logo_path=synk_logo_path, client_logo_path=client_logo_path)

        else:
            # Unknown types render as simple text slide using normalized content
            add_text_slide(prs, meta, sl, synk_logo_path=synk_logo_path, client_logo_path=client_logo_path)

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.read()

def build_base64(deck: dict, filename: str) -> dict:
    data = build_pptx(deck)
    return {
        "filename": filename,
        "file": base64.b64encode(data).decode("utf-8")
    }
